#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
make_presentation.py — builds RapidWatch_RI_Presentation.pptx
An "incredible" dark-ocean-themed deck from the RapidWatch Gulf RI analysis.
Reuses the 7 publication figures in figures/. Run: python make_presentation.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")

# ---------- palette (RapidWatch ocean) ----------
NAVY     = RGBColor(0x0A, 0x18, 0x2B)   # deep background
NAVY2    = RGBColor(0x0E, 0x20, 0x38)   # panel
PANEL    = RGBColor(0x13, 0x2A, 0x45)   # card
PANEL_LT = RGBColor(0x1B, 0x37, 0x57)
INK      = RGBColor(0xEC, 0xF3, 0xFB)   # near-white
MUTE     = RGBColor(0x9C, 0xB6, 0xD2)   # muted blue-grey
FAINT    = RGBColor(0x6B, 0x86, 0xA6)
AMBER    = RGBColor(0xFF, 0xB2, 0x4D)   # RI / warm
AMBER_D  = RGBColor(0xFF, 0x9E, 0x2C)
RED      = RGBColor(0xFF, 0x3B, 0x5C)   # Cat5 / core
TEAL     = RGBColor(0x3C, 0xDD, 0xC9)   # cool / ocean
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

EMU_IN = 914400
SW, SH = 13.333, 7.5
FIG_SIZE = {
    "fig1_hotspot_map.png": (1015, 783),
    "fig2_seasonality.png": (926, 582),
    "fig3_onset_intensity.png": (933, 611),
    "fig4_scaled_d26.png": (1204, 668),
    "fig5_storm_d26.png": (949, 634),
    "fig6_storm_d26_maps.png": (1590, 1241),
    "fig7_zones.png": (933, 807),
}

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color=NAVY):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {})
        sh = ef.makeelement(qn('a:outerShdw'),
                            {'blurRad': '90000', 'dist': '45000', 'dir': '5400000', 'rotWithShape': '0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val': '000000'})
        alpha = clr.makeelement(qn('a:alpha'), {'val': '42000'})
        clr.append(alpha); sh.append(clr); ef.append(sh); el.append(ef)
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (text, size, color, bold, font, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, font, ital) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = ital
            r.font.color.rgb = col; r.font.name = font
    return tb


def R(t, sz, col=INK, bold=False, font="Segoe UI", ital=False):
    return (t, sz, col, bold, font, ital)


def kicker(s, x, y, label, color=AMBER, w=6):
    rect(s, x, y + 0.04, 0.34, 0.06, fill=color)
    txt(s, x + 0.46, y - 0.10, w, 0.4,
        [[R(label.upper(), 13.5, color, True, "Segoe UI Semibold")]])


def footer(s, n):
    rect(s, 0, SH - 0.34, SW, 0.34, fill=NAVY2)
    txt(s, 0.55, SH - 0.40, 8, 0.4,
        [[R("RAPIDWATCH", 9, AMBER, True, "Segoe UI Semibold"),
          R("   ·   Gulf Rapid-Intensification Observatory", 9, FAINT)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, SW - 1.6, SH - 0.40, 1.05, 0.4,
        [[R(f"{n:02d}", 9, MUTE, True)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def fit_image(s, fname, bx, by, bw, bh, frame=True):
    pw, ph = FIG_SIZE[fname]
    ar = pw / ph
    box_ar = bw / bh
    if ar > box_ar:
        w = bw; h = bw / ar
    else:
        h = bh; w = bh * ar
    x = bx + (bw - w) / 2
    y = by + (bh - h) / 2
    if frame:
        pad = 0.12
        rect(s, x - pad, y - pad, w + 2 * pad, h + 2 * pad, fill=WHITE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    s.shapes.add_picture(os.path.join(FIG, fname), Inches(x), Inches(y),
                         Inches(w), Inches(h))
    return x, y, w, h


def stat_card(s, x, y, w, h, value, vcolor, label, sub=None):
    rect(s, x, y, w, h, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, y, 0.09, h, fill=vcolor)
    runs = [[R(value, 33, vcolor, True, "Segoe UI Semibold")],
            [R(label, 12, INK, True)]]
    if sub:
        runs.append([R(sub, 9.5, MUTE)])
    txt(s, x + 0.30, y + 0.18, w - 0.45, h - 0.3, runs, anchor=MSO_ANCHOR.MIDDLE,
        space_after=2, line_spacing=1.0)


def takeaway(s, x, y, w, text, color=TEAL, label="TAKEAWAY"):
    h = 1.15
    rect(s, x, y, w, h, fill=NAVY2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, 0.09, h, fill=color)
    txt(s, x + 0.30, y + 0.14, w - 0.5, h - 0.2,
        [[R(label, 11, color, True, "Segoe UI Semibold")],
         [R(text, 13.5, INK)]], anchor=MSO_ANCHOR.MIDDLE, space_after=4,
        line_spacing=1.05)


def bullets(s, x, y, w, h, items, size=14, gap=9, color=INK, marker=AMBER):
    runs = []
    for it in items:
        runs.append([R("▸  ", size, marker, True), R(it, size, color)])
    txt(s, x, y, w, h, runs, space_after=gap, line_spacing=1.05)


# =========================================================================
# 1 — TITLE
# =========================================================================
s = slide(); bg(s)
# accent ocean band
rect(s, 0, 0, SW, 0.16, fill=AMBER)
rect(s, 0, SH - 0.16, SW, 0.16, fill=PANEL)
txt(s, 0.9, 1.05, 11.5, 0.5,
    [[R("RAPIDWATCH  ·  GULF GEOSPATIAL CANVAS", 15, AMBER, True, "Segoe UI Semibold")]])
txt(s, 0.85, 1.75, 11.7, 2.6,
    [[R("Understanding Hurricane", 50, INK, True, "Segoe UI Semibold")],
     [R("Rapid Intensification", 50, INK, True, "Segoe UI Semibold")]],
    space_after=2, line_spacing=0.98)
rect(s, 0.92, 3.92, 2.6, 0.05, fill=TEAL)
txt(s, 0.9, 4.15, 11.4, 1.2,
    [[R("Where and when hurricanes explode in the Gulf of Mexico — ", 18, MUTE),
      R("and where to point the instruments that could see it coming.", 18, INK, True)]],
    line_spacing=1.15)
# author strip
txt(s, 0.92, 5.7, 11, 0.9,
    [[R("Alex Philp, Ph.D.", 16, INK, True), R("   ·   RapidWatch independent analysis", 16, MUTE)],
     [R("175 years of NOAA best-track data  (1851–2025)  ·  326 Gulf hurricanes", 13, FAINT)]],
    space_after=4)

# =========================================================================
# 2 — THE PROBLEM
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "The problem")
txt(s, 0.55, 0.95, 12, 1.0,
    [[R("The hardest forecast in meteorology", 32, INK, True, "Segoe UI Semibold")]])
txt(s, 0.55, 1.85, 6.4, 4.4,
    [[R("Track forecasts have improved for thirty years. ", 15, MUTE),
      R("Intensity — and rapid intensification above all — has lagged.", 15, INK, True)],
     [R("", 6, MUTE)],
     [R("Rapid intensification turns a modest storm into a major one with little warning. "
        "It is the most dangerous surprise in tropical forecasting — and the Gulf of Mexico "
        "is where it happens most in the United States.", 15, MUTE)],
     [R("", 6, MUTE)],
     [R("Camille ’69 · Katrina & Rita ’05 · Michael ’18 · Ida ’21 · Milton ’24 — "
        "each blew up over the Gulf.", 14, AMBER, True)]],
    space_after=6, line_spacing=1.12)
# right: definition card
rect(s, 7.35, 1.85, 5.45, 3.7, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
rect(s, 7.35, 1.85, 5.45, 0.09, fill=RED)
txt(s, 7.7, 2.15, 4.8, 3.2,
    [[R("WHAT COUNTS AS RAPID INTENSIFICATION", 11.5, RED, True, "Segoe UI Semibold")],
     [R("+30 kt", 44, INK, True, "Segoe UI Semibold")],
     [R("in 24 hours", 19, MUTE)],
     [R("", 6, MUTE)],
     [R("≈ 35 mph of extra wind in a single day — the 95th percentile of all "
        "24-hour changes in the North Atlantic, and the operational standard.", 13, MUTE)]],
    space_after=5, line_spacing=1.1)
footer(s, 2)

# =========================================================================
# 3 — THREE INGREDIENTS
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "The physics", color=TEAL)
txt(s, 0.55, 0.95, 12, 1.0,
    [[R("Three ingredients must align", 32, INK, True, "Segoe UI Semibold")]])
txt(s, 0.55, 1.7, 12, 0.6,
    [[R("The mechanism is well understood. RapidWatch’s question is operational: "
        "where and when do all three coincide?", 15, MUTE)]], line_spacing=1.1)
cards = [
    ("FUEL", AMBER, "Warm sea surface", "Tropical cyclones need water of roughly 26 °C or more — the energy source.", "≥ 26 °C SST"),
    ("AMPLIFIER", RED, "Deep warm water", "Deep warm layers suppress the cold wake that would otherwise cap intensification.", "Ocean heat / D26"),
    ("SWITCH", TEAL, "Low wind shear", "Vertical wind shear is the on/off switch — it tears a developing core apart.", "Weak 200–850 hPa shear"),
]
cx = 0.55; cw = 3.95; gap = 0.27
for i, (tag, col, title, body, foot) in enumerate(cards):
    x = cx + i * (cw + gap)
    rect(s, x, 2.55, cw, 3.45, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, 2.55, cw, 0.10, fill=col)
    txt(s, x + 0.32, 2.85, cw - 0.6, 3.0,
        [[R(tag, 12, col, True, "Segoe UI Semibold")],
         [R(title, 23, INK, True, "Segoe UI Semibold")],
         [R("", 7, MUTE)],
         [R(body, 14, MUTE)]], space_after=5, line_spacing=1.12)
    txt(s, x + 0.32, 5.45, cw - 0.6, 0.5, [[R(foot, 13, col, True)]])
footer(s, 3)

# =========================================================================
# 4 — DATA & METHOD
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "Data & method")
txt(s, 0.55, 0.95, 12, 1.0,
    [[R("Measuring preference, not traffic", 32, INK, True, "Segoe UI Semibold")]])
bullets(s, 0.55, 2.0, 6.4, 4,
    ["175 years of NOAA HURDAT2 best-track data (1851–2025) — 2,004 Atlantic systems.",
     "Restricted to the 326 hurricanes that reached hurricane strength inside a Gulf polygon.",
     "An RI onset = any fix whose wind rises ≥ 30 kt over the next 24 h, placed at the start.",
     "Domain gridded at 0.25°; every cell counts onsets AND all storm fixes passing through.",
     "The quantity is propensity — onsets ÷ traffic — so a busy shipping lane can’t fake a hotspot."],
    size=14.5, gap=11)
# right method callout
rect(s, 7.4, 2.0, 5.4, 3.6, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
rect(s, 7.4, 2.0, 0.09, 3.6, fill=TEAL)
txt(s, 7.75, 2.3, 4.85, 3.1,
    [[R("WHY NORMALIZE BY TRAFFIC?", 12, TEAL, True, "Segoe UI Semibold")],
     [R("", 5, MUTE)],
     [R("Storms cluster along certain paths. Counting raw onsets would just rediscover "
        "where storms travel.", 14, MUTE)],
     [R("", 5, MUTE)],
     [R("Dividing by the number of storms that passed through each cell isolates a real "
        "preference — the places where storms are unusually ", 14, MUTE),
      R("likely to ignite", 14, AMBER, True),
      R(", not merely present.", 14, MUTE)]],
    space_after=4, line_spacing=1.12)
footer(s, 4)

# =========================================================================
# 5 — HEADLINE NUMBERS
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "Headline numbers", color=AMBER)
txt(s, 0.55, 0.95, 12, 1.0,
    [[R("What 175 years of storms reveal", 32, INK, True, "Segoe UI Semibold")]])
data = [
    ("3.1×", RED,   "More likely in the core", "RI onset in the southern-Gulf core vs. basin average"),
    ("47%", AMBER,  "Of Gulf hurricanes", "underwent at least one RI onset inside the Gulf"),
    ("84%", TEAL,   "Of onsets in Aug–Oct", "peaking sharply in September"),
    ("50 kt", AMBER,"Median strength at onset", "RI begins while a storm is still a tropical storm"),
    ("66%", RED,    "Begin at ≤ TS strength", "the storms that explode rarely look dangerous yet"),
    ("326", TEAL,   "Gulf hurricanes analyzed", "411 RI-onset intervals detected, 1851–2025"),
]
gx, gy = 0.55, 2.0
cw, ch = 3.95, 1.75; gxp, gyp = 0.27, 0.27
for i, (v, c, lab, sub) in enumerate(data):
    r, col = divmod(i, 3)
    x = gx + col * (cw + gxp); y = gy + r * (ch + gyp)
    stat_card(s, x, y, cw, ch, v, c, lab, sub)
footer(s, 5)

# =========================================================================
# 6 — FINDING 1: WHERE  (fig1)
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "Finding 1 — Where", color=RED)
txt(s, 0.55, 0.95, 7, 1.0,
    [[R("It is geographically concentrated", 30, INK, True, "Segoe UI Semibold")]])
txt(s, 0.55, 1.85, 5.0, 3.0,
    [[R("Nearly ", 15, MUTE), R("one fix in four", 15, RED, True),
      R(" in the southern-Gulf core is an RI onset — ", 15, MUTE),
      R("3.1× the basin average.", 15, INK, True)],
     [R("", 6, MUTE)],
     [R("The tightest core sits over the Bay of Campeche and the southwestern flank of "
        "the Loop Current (≈ 21–23 °N, 92–94 °W), fanning east toward the Yucatán Channel.", 14, MUTE)]],
    space_after=6, line_spacing=1.13)
takeaway(s, 0.55, 5.05, 5.0,
         "A genuine preference — it survives normalization by storm traffic.", color=RED)
fit_image(s, "fig1_hotspot_map.png", 5.95, 1.55, 6.95, 4.95)
footer(s, 6)

# =========================================================================
# 7 — FINDING 2: WHEN  (fig2)
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "Finding 2 — When", color=AMBER)
txt(s, 0.55, 0.95, 7, 1.0,
    [[R("It is strongly seasonal", 30, INK, True, "Segoe UI Semibold")]])
txt(s, 0.55, 1.85, 5.0, 3.0,
    [[R("September is the peak month; ", 15, MUTE),
      R("84% of all onsets fall in August–October.", 15, INK, True)],
     [R("", 6, MUTE)],
     [R("This is when the three ingredients align at once: sea-surface temperature at its "
        "annual peak, the Loop Current carrying its deepest warm layer, and basin-wide shear "
        "at its climatological lowest.", 14, MUTE)]],
    space_after=6, line_spacing=1.13)
takeaway(s, 0.55, 5.05, 5.0,
         "One calendar: pre-position for the warm season, peaking in September.", color=AMBER)
fit_image(s, "fig2_seasonality.png", 5.95, 1.7, 6.95, 4.6)
footer(s, 7)

# =========================================================================
# 8 — FINDING 3: AT WHAT INTENSITY  (fig3)
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "Finding 3 — At what strength", color=TEAL)
txt(s, 0.55, 0.95, 7.2, 1.0,
    [[R("It starts in weak systems", 30, INK, True, "Segoe UI Semibold")]])
txt(s, 0.55, 1.85, 5.0, 3.0,
    [[R("Median intensity at onset is just ", 15, MUTE), R("50 kt", 15, TEAL, True),
      R(". Two-thirds of onsets begin at tropical-storm strength or below; 87% at Cat 1 or below.", 15, MUTE)],
     [R("", 6, MUTE)],
     [R("By the time a storm looks dangerous, the window most in need of observation has "
        "often already opened.", 14, INK, True)]],
    space_after=6, line_spacing=1.13)
takeaway(s, 0.55, 5.05, 5.0,
         "Targeting must include developing systems — not just the storms that already look like a threat.",
         color=TEAL)
fit_image(s, "fig3_onset_intensity.png", 5.95, 1.7, 6.95, 4.6)
footer(s, 8)

# =========================================================================
# 9 — FINDING 4: WHY  (fig6 -> "Figure 4")
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "Finding 4 — Why", color=RED)
txt(s, 0.55, 0.95, 8, 1.0,
    [[R("The hotspot sits on the deep warm water", 28, INK, True, "Segoe UI Semibold")]])
txt(s, 0.55, 1.8, 4.55, 3.6,
    [[R("The geography is no coincidence. The core overlies the Gulf’s two deepest "
        "warm-water features — the ", 14, MUTE), R("Loop Current", 14, RED, True),
      R(" and the ", 14, MUTE), R("Bay of Campeche warm pool.", 14, RED, True)],
     [R("", 5, MUTE)],
     [R("Measured ocean heat (depth of the 26 °C isotherm, D26) confirms it storm by storm: "
        "Katrina, Rita, and Helene each intensified over water 1.2–1.7× the Gulf-median depth.", 13.5, MUTE)]],
    space_after=5, line_spacing=1.12)
fit_image(s, "fig6_storm_d26_maps.png", 5.25, 1.5, 7.7, 5.4)
footer(s, 9)

# =========================================================================
# 10 — THE MILTON EXCEPTION (fig5)
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "The instructive exception", color=AMBER)
txt(s, 0.55, 0.95, 8, 1.0,
    [[R("Milton 2024 — deep heat is not the only trigger", 26, INK, True, "Segoe UI Semibold")]])
txt(s, 0.55, 1.85, 5.0, 3.4,
    [[R("The most explosive case on record — ", 15, MUTE), R("+80 kt in a day", 15, AMBER, True),
      R(" — began over the Bay of Campeche where the warm layer was only ~52 m deep, "
        "below the Gulf median.", 15, MUTE)],
     [R("", 6, MUTE)],
     [R("Its fuel was record sea-surface temperature and near-zero shear — not the "
        "deepest ocean heat.", 14, INK, True)]],
    space_after=6, line_spacing=1.13)
takeaway(s, 0.55, 5.2, 5.0,
         "Deep warm water amplifies rapid intensification — but it is not the sole trigger.",
         color=AMBER)
fit_image(s, "fig5_storm_d26.png", 5.95, 1.7, 6.95, 4.6)
footer(s, 10)

# =========================================================================
# 11 — THE POPULATION TEST (fig4 -> "Figure 6")
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "The honest test", color=TEAL)
txt(s, 0.55, 0.95, 8.6, 1.0,
    [[R("Does warm-water depth add forecasting skill?", 26, INK, True, "Segoe UI Semibold")]])
txt(s, 0.55, 1.8, 4.55, 4.2,
    [[R("Across 76 Gulf hurricanes (698 fixes, 89 onsets), depth alone is ", 14, MUTE),
      R("barely better than a coin flip", 14, TEAL, True),
      R(" at separating onsets from non-onsets (AUC 0.55).", 14, MUTE)],
     [R("", 5, MUTE)],
     [R("Add D26 to location, season, and current intensity, and out-of-sample skill is "
        "essentially unchanged (0.693 → 0.689).", 14, MUTE)],
     [R("", 5, MUTE)],
     [R("Not irrelevance — redundancy. Where and when already encode where the deep warm "
        "water lies.", 14, INK, True)]],
    space_after=5, line_spacing=1.12)
fit_image(s, "fig4_scaled_d26.png", 5.25, 1.75, 7.7, 4.6)
footer(s, 11)

# =========================================================================
# 12 — TWO OPERATIONAL ZONES (fig7)
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "From science to strategy", color=AMBER)
txt(s, 0.55, 0.95, 8, 1.0,
    [[R("Two distinct operational zones", 30, INK, True, "Segoe UI Semibold")]])
# zone A card
rect(s, 0.55, 1.85, 5.0, 1.85, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
rect(s, 0.55, 1.85, 0.09, 1.85, fill=RED)
txt(s, 0.85, 2.02, 4.5, 1.6,
    [[R("ZONE A — ONSET", 12, RED, True, "Segoe UI Semibold")],
     [R("Southern-Gulf core", 18, INK, True, "Segoe UI Semibold")],
     [R("Where RI begins, in weak systems entering from the south. Sample the ocean "
        "ahead of the storm to catch it early.", 12.5, MUTE)]],
    space_after=3, line_spacing=1.08)
# zone B card
rect(s, 0.55, 3.85, 5.0, 1.85, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
rect(s, 0.55, 3.85, 0.09, 1.85, fill=AMBER)
txt(s, 0.85, 4.02, 4.5, 1.6,
    [[R("ZONE B — IMPACT", 12, AMBER, True, "Segoe UI Semibold")],
     [R("Loop Current corridor & shelf", 17, INK, True, "Segoe UI Semibold")],
     [R("Where RI continues toward U.S. landfall — Katrina, Rita, Michael, Ida. "
        "Highest stakes for warning.", 12.5, MUTE)]],
    space_after=3, line_spacing=1.08)
fit_image(s, "fig7_zones.png", 5.95, 1.55, 6.95, 5.0)
footer(s, 12)

# =========================================================================
# 13 — THE PROTOCOL
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "Recommendation", color=AMBER)
txt(s, 0.55, 0.95, 12, 1.0,
    [[R("The Gulf RI Watch protocol", 32, INK, True, "Segoe UI Semibold")]])
cols = [
    ("WHEN — TRIGGER", TEAL,
     ["Any system forms in the Gulf or Bay of Campeche — any intensity, including disturbances.",
      "A NW-Caribbean system is forecast through the Yucatán Channel within 72 h.",
      "A Gulf system is forecast over the Loop Current toward the north-central shelf."]),
    ("WHERE — TARGET", RED,
     ["Primary: the southern-Gulf core box (20.9–22.9 °N, 91.6–93.9 °W).",
      "Expanded to the watch box to take in the Yucatán intrusion.",
      "Secondary: the Loop Current corridor along the forecast track."]),
    ("HOW — LEAD & ORDER", AMBER,
     ["Sample 12–36 h ahead of projected onset, ahead of the track.",
      "In warm season, treat any qualifying system as RI-capable by default.",
      "Priority: ocean heat → inner core of weak systems → environmental shear."]),
]
cx = 0.55; cw = 3.95; gp = 0.27
for i, (tag, col, items) in enumerate(cols):
    x = cx + i * (cw + gp)
    rect(s, x, 2.0, cw, 4.4, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, 2.0, cw, 0.10, fill=col)
    txt(s, x + 0.3, 2.32, cw - 0.55, 0.5, [[R(tag, 12.5, col, True, "Segoe UI Semibold")]])
    runs = [[R("▸  ", 12.5, col, True), R(it, 12.5, INK)] for it in items]
    txt(s, x + 0.3, 2.95, cw - 0.55, 3.3, runs, space_after=10, line_spacing=1.08)
footer(s, 13)

# =========================================================================
# 14 — OBSERVATIONS NEEDED
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "Closing the gaps", color=TEAL)
txt(s, 0.55, 0.95, 12, 1.0,
    [[R("The platforms — by the two persistent gaps", 30, INK, True, "Segoe UI Semibold")]])
# Gap 1
rect(s, 0.55, 1.95, 6.05, 4.4, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
rect(s, 0.55, 1.95, 6.05, 0.10, fill=TEAL)
txt(s, 0.9, 2.25, 5.4, 0.6,
    [[R("GAP 1", 12, TEAL, True, "Segoe UI Semibold"),
      R("   Subsurface ocean heat ahead of the storm", 14.5, INK, True)]])
bullets(s, 0.9, 2.95, 5.4, 3.2,
    ["AXBT / AXCTD ocean profilers dropped from P-3 & C-130 across the core box",
     "Air-deployed ALAMO / APEX profiling floats through the event",
     "IOOS hurricane-glider picket line, Yucatán to Campeche Bank",
     "Satellite altimetry — Sentinel-6, Jason, SWOT — for ocean-heat mapping",
     "Enhanced seasonal Argo to anchor altimetry → warm-layer depth"],
    size=13, gap=9, marker=TEAL)
# Gap 2
rect(s, 6.85, 1.95, 5.95, 4.4, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
rect(s, 6.85, 1.95, 5.95, 0.10, fill=AMBER)
txt(s, 7.2, 2.25, 5.3, 0.6,
    [[R("GAP 2", 12, AMBER, True, "Segoe UI Semibold"),
      R("   Inner core of developing systems", 14.5, INK, True)]])
bullets(s, 7.2, 2.95, 5.3, 3.2,
    ["NOAA WP-3D with Tail Doppler Radar + dropsondes — tasked earlier, at TS strength",
     "Small uncrewed aircraft (Altius-600, Coyote) in the boundary-layer inflow",
     "Saildrone surface vehicles for air–sea flux through the season",
     "G-IV synoptic surveillance to constrain shear and steering",
     "Protect passive-microwave coverage (DMSP/SSMIS ending) over the Gulf"],
    size=13, gap=9, marker=AMBER)
footer(s, 14)

# =========================================================================
# 15 — CONCLUSION / BOTTOM LINE
# =========================================================================
s = slide(); bg(s); kicker(s, 0.55, 0.55, "Bottom line")
txt(s, 0.55, 0.95, 12.2, 1.0,
    [[R("A seasonal calendar and a priority map", 32, INK, True, "Segoe UI Semibold")]])
points = [
    ("Concentrated", RED, "RI onset is ~3× more likely in a southern-Gulf core (≈ 21–23 °N, 92–94 °W)."),
    ("Seasonal", AMBER, "84% of onsets fall Aug–Oct, peaking in September."),
    ("Early", TEAL, "It begins in weak systems — two-thirds at tropical-storm strength or below."),
    ("Redundant signal", AMBER, "Measured ocean heat adds little skill beyond location + season — they already encode it."),
]
gy = 2.0
for i, (h, c, b) in enumerate(points):
    y = gy + i * 1.02
    rect(s, 0.55, y, 7.3, 0.88, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 0.55, y, 0.09, 0.88, fill=c)
    txt(s, 0.85, y, 7.0, 0.88,
        [[R(h + "  —  ", 14.5, c, True), R(b, 13.5, INK)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
# right summary panel
rect(s, 8.15, 2.0, 4.65, 4.06, fill=NAVY2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
rect(s, 8.15, 2.0, 4.65, 0.10, fill=TEAL)
txt(s, 8.5, 2.35, 4.0, 3.5,
    [[R("THE OPERATIONAL PAYOFF", 11.5, TEAL, True, "Segoe UI Semibold")],
     [R("", 6, MUTE)],
     [R("Point finite flight-hours and glider-days at the high-propensity core, in the "
        "right season, at 12–36 h lead.", 14, INK)],
     [R("", 6, MUTE)],
     [R("The value of real-time ocean heat is in ", 13.5, MUTE),
      R("initializing the coupled forecast models", 13.5, AMBER, True),
      R(" — not as a stand-alone warning signal.", 13.5, MUTE)]],
    space_after=4, line_spacing=1.13)
footer(s, 15)

# =========================================================================
# 16 — CLOSING
# =========================================================================
s = slide(); bg(s)
rect(s, 0, 0, SW, 0.16, fill=AMBER)
txt(s, 0.9, 2.5, 11.5, 2.0,
    [[R("RapidWatch", 46, INK, True, "Segoe UI Semibold")],
     [R("Gulf Geospatial Canvas — Rapid-Intensification Observatory", 18, MUTE)]],
    space_after=8)
rect(s, 0.95, 4.15, 2.6, 0.05, fill=TEAL)
txt(s, 0.92, 4.45, 11.4, 1.6,
    [[R("Interactive map · 326-storm history · live buoys & NHC cone · four measured ocean-heat fields",
        14.5, MUTE)],
     [R("aphilp1.github.io/rapidwatch", 15, AMBER, True)],
     [R("", 6, MUTE)],
     [R("Alex Philp, Ph.D.  ·  Data: NOAA HURDAT2 1851–2025 · HYCOM GOFS 3.1 · Copernicus GLORYS",
        12, FAINT)]],
    space_after=6, line_spacing=1.2)

out = os.path.join(HERE, "RapidWatch_RI_Presentation.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
